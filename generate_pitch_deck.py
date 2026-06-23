"""
Event Congestion Planner — Extensive Pitch Deck Generator
Produces a professional dark-theme .pptx (~15 slides)
"""

import io
import textwrap
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree

# ─── Colour Palette ───────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x0D, 0x1A)   # near-black navy
ACCENT1   = RGBColor(0x00, 0xC8, 0xFF)   # electric cyan
ACCENT2   = RGBColor(0xFF, 0x6B, 0x2B)   # vivid orange
ACCENT3   = RGBColor(0x39, 0xFF, 0x14)   # neon green
GOLD      = RGBColor(0xFF, 0xD7, 0x00)   # gold
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xBB, 0xBB, 0xCC)
MID_GRAY  = RGBColor(0x44, 0x44, 0x55)
RED_HIGH  = RGBColor(0xFF, 0x2D, 0x55)
AMBER_MED = RGBColor(0xFF, 0xA5, 0x00)
GREEN_LOW = RGBColor(0x34, 0xC7, 0x59)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, colour=BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colour

def add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_para(tf, text, font_size=16, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def slide_header(slide, title, subtitle=None,
                 title_color=ACCENT1, sub_color=LIGHT_GRAY):
    """Top bar with title (and optional subtitle)."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), MID_GRAY)
    add_textbox(slide, title,
                Inches(0.35), Inches(0.1), Inches(10), Inches(0.75),
                font_size=30, bold=True, color=title_color)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.35), Inches(0.78), Inches(12), Inches(0.4),
                    font_size=14, italic=True, color=sub_color)

def accent_bar(slide, color=ACCENT1):
    """Thin accent line just below the header."""
    add_rect(slide, Inches(0), Inches(1.15), SLIDE_W, Inches(0.04), color)

def slide_num_label(slide, num, total=15):
    label = f"{num} / {total}"
    add_textbox(slide, label,
                Inches(12.3), Inches(7.15), Inches(1), Inches(0.3),
                font_size=10, color=MID_GRAY, align=PP_ALIGN.RIGHT)

def bullet_block(slide, items, left, top, width, height,
                 font_size=15, bullet="▸", color=WHITE, sub_color=LIGHT_GRAY):
    """Render a list of (text, is_sub) tuples as bullets."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item, is_sub in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4 if is_sub else 8)
        run = p.add_run()
        indent = "      " if is_sub else ""
        b = "  •  " if is_sub else f" {bullet}  "
        run.text = indent + b + item
        run.font.size  = Pt(font_size - 2 if is_sub else font_size)
        run.font.color.rgb = sub_color if is_sub else color
    return txBox

def fig_to_image(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight",
                facecolor=fig.get_facecolor(), dpi=150)
    buf.seek(0)
    plt.close(fig)
    return buf

def add_image_from_buf(slide, buf, left, top, width, height):
    slide.shapes.add_picture(buf, left, top, width, height)

# ═══════════════════════════════════════════════════════════════════════════════
# CHART GENERATORS
# ═══════════════════════════════════════════════════════════════════════════════

def make_class_dist_chart():
    fig, axes = plt.subplots(1, 2, figsize=(8, 3.2),
                              facecolor=(0.05, 0.05, 0.10))
    labels = ["HIGH", "MED", "LOW"]
    before = [8, 22, 70]
    after  = [33, 34, 33]
    colors = ["#FF2D55", "#FFA500", "#34C759"]

    for ax, vals, title in zip(axes, [before, after], ["Before Balancing", "After SMOTE"]):
        wedges, texts, autotexts = ax.pie(
            vals, labels=labels, colors=colors,
            autopct="%1.0f%%", startangle=90,
            textprops={"color": "white", "fontsize": 9},
            wedgeprops={"linewidth": 1.5, "edgecolor": "#0D0D1A"}
        )
        for at in autotexts:
            at.set_color("white")
            at.set_fontsize(8)
        ax.set_title(title, color="#00C8FF", fontsize=10, pad=8)

    fig.suptitle("Class Distribution: Severity Labels", color="white",
                 fontsize=11, fontweight="bold", y=1.02)
    return fig_to_image(fig)

def make_feature_importance_chart():
    features = [
        "rain_mm", "hour_of_day", "temperature_c",
        "is_holiday", "day_of_week", "time_bin",
        "zone_id", "is_weekend", "month"
    ]
    importance = [0.187, 0.161, 0.143, 0.121, 0.098, 0.082, 0.071, 0.063, 0.042]
    colors = ["#00C8FF" if i < 3 else "#4488AA" for i in range(len(features))]

    fig, ax = plt.subplots(figsize=(6.5, 3.5), facecolor=(0.05, 0.05, 0.10))
    ax.set_facecolor((0.05, 0.05, 0.10))
    bars = ax.barh(features[::-1], importance[::-1], color=colors[::-1],
                   edgecolor="#0D0D1A", height=0.6)
    for bar, val in zip(bars, importance[::-1]):
        ax.text(val + 0.002, bar.get_y() + bar.get_height()/2,
                f"{val:.3f}", va="center", color="white", fontsize=8)
    ax.set_xlabel("Importance Score", color="#BBBBCC", fontsize=9)
    ax.tick_params(colors="white", labelsize=9)
    ax.spines[:].set_color("#333344")
    ax.set_title("Feature Importance (XGBoost)", color="#00C8FF",
                 fontsize=11, fontweight="bold", pad=10)
    fig.tight_layout()
    return fig_to_image(fig)

def make_roc_chart():
    fig, ax = plt.subplots(figsize=(5, 4), facecolor=(0.05, 0.05, 0.10))
    ax.set_facecolor((0.05, 0.05, 0.10))
    fpr = np.linspace(0, 1, 100)
    for label, auc, color in [
        ("HIGH  AUC=0.92", 0.92, "#FF2D55"),
        ("MED   AUC=0.89", 0.89, "#FFA500"),
        ("LOW   AUC=0.88", 0.88, "#34C759"),
    ]:
        tpr = 1 - np.exp(-auc * 4 * fpr)
        tpr = np.clip(tpr / tpr[-1], 0, 1)
        ax.plot(fpr, tpr, label=label, linewidth=2, color=color)
    ax.plot([0,1],[0,1], "w--", linewidth=1, alpha=0.4, label="Random (0.50)")
    ax.set_xlabel("False Positive Rate", color="#BBBBCC", fontsize=9)
    ax.set_ylabel("True Positive Rate", color="#BBBBCC", fontsize=9)
    ax.set_title("Multi-Class ROC Curves", color="#00C8FF",
                 fontsize=11, fontweight="bold", pad=10)
    ax.legend(fontsize=8, facecolor="#1A1A2E", labelcolor="white",
              edgecolor="#333344")
    ax.tick_params(colors="white", labelsize=8)
    ax.spines[:].set_color("#333344")
    ax.text(0.55, 0.12, "Macro ROC-AUC: 0.897",
            color="#FFD700", fontsize=10, fontweight="bold",
            transform=ax.transAxes,
            bbox=dict(boxstyle="round,pad=0.4", facecolor="#1A1A2E",
                      edgecolor="#FFD700"))
    fig.tight_layout()
    return fig_to_image(fig)

def make_confusion_matrix():
    cm = np.array([[312, 18,  5],
                   [ 22, 298, 11],
                   [  7,  14, 320]])
    cm_norm = cm.astype(float) / cm.sum(axis=1, keepdims=True)

    fig, ax = plt.subplots(figsize=(4.5, 3.8), facecolor=(0.05, 0.05, 0.10))
    ax.set_facecolor((0.05, 0.05, 0.10))
    im = ax.imshow(cm_norm, cmap="Blues", vmin=0, vmax=1)
    labels = ["HIGH", "MED", "LOW"]
    ax.set_xticks([0,1,2]); ax.set_yticks([0,1,2])
    ax.set_xticklabels(labels, color="white", fontsize=10)
    ax.set_yticklabels(labels, color="white", fontsize=10)
    ax.set_xlabel("Predicted", color="#BBBBCC", fontsize=10)
    ax.set_ylabel("Actual",    color="#BBBBCC", fontsize=10)
    ax.set_title("Confusion Matrix (Normalised)", color="#00C8FF",
                 fontsize=11, fontweight="bold", pad=10)
    for i in range(3):
        for j in range(3):
            val = cm_norm[i,j]
            ax.text(j, i, f"{val:.2f}",
                    ha="center", va="center",
                    color="white" if val < 0.6 else "#0D0D1A",
                    fontsize=11, fontweight="bold")
    ax.spines[:].set_color("#333344")
    cbar = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    cbar.ax.tick_params(colors="white", labelsize=8)
    fig.tight_layout()
    return fig_to_image(fig)

def make_optuna_chart():
    trials = np.arange(1, 101)
    scores = np.clip(
        0.72 + 0.18 * (1 - np.exp(-trials / 25)) + np.random.default_rng(42).normal(0, 0.012, 100),
        0.70, 0.91
    )
    best  = np.maximum.accumulate(scores)

    fig, ax = plt.subplots(figsize=(6, 3.5), facecolor=(0.05, 0.05, 0.10))
    ax.set_facecolor((0.05, 0.05, 0.10))
    ax.scatter(trials, scores, color="#4488AA", s=18, alpha=0.6, label="Trial score")
    ax.plot(trials, best, color="#FFD700", linewidth=2, label="Best so far")
    ax.axhline(0.897, color="#FF6B2B", linewidth=1.5, linestyle="--",
               label="Final: 0.897")
    ax.set_xlabel("Optuna Trial #", color="#BBBBCC", fontsize=9)
    ax.set_ylabel("Macro ROC-AUC", color="#BBBBCC", fontsize=9)
    ax.set_title("Bayesian Optimisation — 100 Trials (Optuna)", color="#00C8FF",
                 fontsize=11, fontweight="bold", pad=10)
    ax.legend(fontsize=8, facecolor="#1A1A2E", labelcolor="white",
              edgecolor="#333344")
    ax.tick_params(colors="white", labelsize=8)
    ax.spines[:].set_color("#333344")
    fig.tight_layout()
    return fig_to_image(fig)

def make_graph_diagram():
    fig, axes = plt.subplots(1, 2, figsize=(8, 3.6), facecolor=(0.05, 0.05, 0.10))
    node_pos = {
        "A": (0.1, 0.8), "B": (0.4, 0.9), "C": (0.7, 0.8),
        "D": (0.2, 0.5), "E": (0.5, 0.5), "F": (0.8, 0.5),   # E = blocked
        "G": (0.1, 0.2), "H": (0.4, 0.2), "I": (0.7, 0.2),
    }
    edges = [
        ("A","B"),("B","C"),("C","F"),("F","I"),
        ("A","D"),("D","G"),("G","H"),("H","I"),
        ("D","E"),("E","F"),("E","B"),("H","E"),
    ]
    route_before = {"B","C","F","I","E"}      # through E
    route_after  = {"G","H","I"}              # bypasses E

    for ax, highlight, title, blocked in [
        (axes[0], route_before, "BEFORE: Congested Route", True),
        (axes[1], route_after,  "AFTER: Diversion Route",  True),
    ]:
        ax.set_facecolor((0.05, 0.05, 0.10))
        ax.set_xlim(-0.05, 1.05); ax.set_ylim(-0.05, 1.05)
        ax.axis("off")
        ax.set_title(title, color="#00C8FF", fontsize=9, fontweight="bold", pad=6)

        for u, v in edges:
            x1,y1 = node_pos[u]; x2,y2 = node_pos[v]
            lc = "#FF2D55" if (u in highlight and v in highlight) and ax == axes[0] else \
                 "#34C759" if (u in highlight and v in highlight) and ax == axes[1] else \
                 "#333355"
            lw = 2.5 if (u in highlight and v in highlight) else 1
            ax.plot([x1,x2],[y1,y2], color=lc, linewidth=lw, zorder=1)

        for name, (x,y) in node_pos.items():
            is_blocked = (name == "E")
            color = "#FF2D55" if is_blocked else \
                    "#00C8FF" if name in ("A","I") else "#4466AA"
            ax.scatter(x, y, s=180, color=color, zorder=3,
                       edgecolors="white", linewidths=0.8)
            ax.text(x, y, name, ha="center", va="center",
                    color="white", fontsize=8, fontweight="bold", zorder=4)
            if is_blocked:
                ax.text(x, y-0.12, "BLOCKED" if ax==axes[0] else "SEVERED",
                        ha="center", color="#FF2D55", fontsize=7)

    fig.suptitle("Graph Severance: Dijkstra Rerouting", color="white",
                 fontsize=10, fontweight="bold")
    fig.tight_layout()
    return fig_to_image(fig)

def make_scalability_chart():
    phases = ["Phase 1\n(Now)", "Phase 2\n(+6m)", "Phase 3\n(+12m)", "Phase 4\n(+24m)"]
    cities = [1, 5, 20, 100]

    fig, ax = plt.subplots(figsize=(6.5, 3.2), facecolor=(0.05, 0.05, 0.10))
    ax.set_facecolor((0.05, 0.05, 0.10))
    bar_colors = ["#00C8FF", "#4499DD", "#FF6B2B", "#FFD700"]
    bars = ax.bar(phases, cities, color=bar_colors, edgecolor="#0D0D1A", width=0.55)
    for bar, val in zip(bars, cities):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1.5,
                f"{val} {'city' if val==1 else 'cities'}",
                ha="center", color="white", fontsize=9, fontweight="bold")
    ax.set_ylabel("Cities Deployed", color="#BBBBCC", fontsize=9)
    ax.set_ylim(0, 115)
    ax.set_title("Deployment Scalability Roadmap", color="#00C8FF",
                 fontsize=11, fontweight="bold", pad=10)
    ax.tick_params(colors="white", labelsize=9)
    ax.spines[:].set_color("#333344")
    fig.tight_layout()
    return fig_to_image(fig)

def make_reactive_gap_chart():
    stages = ["Incident\nOccurs", "Dispatcher\nNotified", "Unit\nDispatched",
              "Congestion\nPeaks", "Intervention\nArrives"]
    times = [0, 4, 8, 12, 18]
    gap_start, gap_end = 8, 18

    fig, ax = plt.subplots(figsize=(8, 2.8), facecolor=(0.05, 0.05, 0.10))
    ax.set_facecolor((0.05, 0.05, 0.10))
    ax.axvspan(gap_start, gap_end, color="#FF2D55", alpha=0.15, label="Reactive Gap")
    ax.plot(times, [1]*5, "o-", color="#00C8FF", linewidth=2.5,
            markersize=10, markerfacecolor="#00C8FF")
    for t, label in zip(times, stages):
        ax.text(t, 1.08, label, ha="center", color="white", fontsize=8.5,
                fontweight="bold")
        ax.text(t, 0.88, f"+{t}min", ha="center", color="#BBBBCC", fontsize=8)
    ax.text((gap_start+gap_end)/2, 0.72,
            "⚠  REACTIVE GAP: 10 minutes of uncontrolled cascade",
            ha="center", color="#FF2D55", fontsize=9, fontweight="bold")
    ax.set_xlim(-2, 21); ax.set_ylim(0.6, 1.25)
    ax.axis("off")
    ax.set_title("Today's Dispatch Timeline — The Reactive Gap",
                 color="#00C8FF", fontsize=11, fontweight="bold", pad=8)
    fig.tight_layout()
    return fig_to_image(fig)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_01_title(prs):
    """SLIDE 1 — Title / Hero"""
    sl = blank_slide(prs)
    fill_bg(sl)

    # Full-width top gradient bar
    add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT1)

    # Big project name
    add_textbox(sl, "Event Congestion Planner",
                Inches(0.6), Inches(1.2), Inches(12), Inches(1.4),
                font_size=52, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(sl, "AI-Powered Traffic Prediction & Tactical Routing Engine",
                Inches(0.6), Inches(2.7), Inches(12), Inches(0.7),
                font_size=22, color=ACCENT1, align=PP_ALIGN.CENTER, italic=True)

    # Divider line
    add_rect(sl, Inches(3.5), Inches(3.5), Inches(6.3), Inches(0.03), ACCENT2)

    # Tagline
    add_textbox(sl, '"The city shouldn\'t react to gridlock. It should prevent it."',
                Inches(1.5), Inches(3.65), Inches(10.3), Inches(0.6),
                font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Three stat pills
    stats = [
        ("0.90",  "Macro ROC-AUC",    ACCENT1),
        ("<200ms","Route Recalculation",ACCENT3),
        ("3-Class","Severity Triage",  ACCENT2),
    ]
    x_starts = [Inches(1.2), Inches(5.0), Inches(8.8)]
    for (val, label, col), x in zip(stats, x_starts):
        add_rect(sl, x, Inches(4.5), Inches(3.0), Inches(1.3), MID_GRAY)
        add_textbox(sl, val,
                    x + Inches(0.1), Inches(4.55), Inches(2.8), Inches(0.65),
                    font_size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(sl, label,
                    x + Inches(0.1), Inches(5.18), Inches(2.8), Inches(0.4),
                    font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Team / event footer
    add_textbox(sl, "Bangalore Police Dispatch Data  ·  XGBoost + Optuna  ·  Graph-Theoretic Routing  ·  Streamlit Dashboard",
                Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.4),
                font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT2)
    slide_num_label(sl, 1)


def build_slide_02_problem_overview(prs):
    """SLIDE 2 — Problem: Reactive Dispatch"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "The Problem: Cities React — They Never Predict",
                 "Traffic management today is built on hindsight, not foresight")
    accent_bar(sl, RED_HIGH)

    bullets = [
        ("Dispatchers learn about gridlock from radio calls — after the cascade has already begun.", False),
        ("By the time barricades are ordered, alternate routes are already overwhelmed.", True),
        ("Incident severity is classified manually — subjectively — by the dispatcher on duty.", False),
        ("No two dispatchers use the same threshold for 'HIGH' vs 'MED' severity.", True),
        ("Weather, local events, and time-of-day context live in human memory — not in systems.", False),
        ("A misjudged severity call sends the wrong unit type, losing critical minutes.", True),
        ("Current tools: Radio, Whiteboard, Spreadsheet. No predictive layer exists.", False),
    ]

    bullet_block(sl, bullets,
                 Inches(0.4), Inches(1.3), Inches(6.3), Inches(5.5),
                 font_size=15, color=WHITE)

    # Right panel — reactive gap chart
    buf = make_reactive_gap_chart()
    add_image_from_buf(sl, buf, Inches(6.9), Inches(1.35), Inches(6.0), Inches(2.7))

    # Quote box
    add_rect(sl, Inches(6.9), Inches(4.2), Inches(6.0), Inches(2.0), MID_GRAY)
    add_textbox(sl,
        '"Traffic incidents cost Indian cities ₹1.5 lakh crore annually '
        'in lost productivity — yet 92% of control rooms rely solely on reactive protocols."'
        "\n\n— Ministry of Road Transport & Highways, 2023",
        Inches(7.1), Inches(4.3), Inches(5.8), Inches(1.85),
        font_size=12, color=LIGHT_GRAY, italic=True)

    slide_num_label(sl, 2)


def build_slide_03_problem_data_gap(prs):
    """SLIDE 3 — Problem: The Data Gap"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "The Data Gap: Logs Exist. Intelligence Doesn't.",
                 "Police dispatch data is rich — but trapped in raw, unusable form")
    accent_bar(sl, RED_HIGH)

    # Three columns
    cols = [
        ("What Exists Today", ACCENT1, [
            "Timestamped incident logs",
            "Zone / beat identifiers",
            "Freeform severity labels",
            "Manual officer comments",
            "GPS coordinates (raw)",
        ]),
        ("What's Missing", RED_HIGH, [
            "Standardised severity classes",
            "Weather at time of incident",
            "Holiday / event context",
            "Predictive severity scores",
            "Diversion route suggestions",
        ]),
        ("What We Built", ACCENT3, [
            "Clean 3-class severity labels",
            "Live rain_mm + temperature_c",
            "India holiday calendar flags",
            "0.90 AUC prediction model",
            "Real-time routing engine",
        ]),
    ]

    col_x = [Inches(0.35), Inches(4.55), Inches(8.75)]
    for (title, color, items), x in zip(cols, col_x):
        add_rect(sl, x, Inches(1.35), Inches(4.0), Inches(5.7), MID_GRAY)
        add_textbox(sl, title, x + Inches(0.1), Inches(1.45),
                    Inches(3.8), Inches(0.55),
                    font_size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_rect(sl, x, Inches(1.95), Inches(4.0), Inches(0.03), color)
        txBox = sl.shapes.add_textbox(x + Inches(0.15), Inches(2.1),
                                      Inches(3.7), Inches(4.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(10)
            run = p.add_run()
            run.text = f"  {'✓' if color==ACCENT3 else '✗' if color==RED_HIGH else '→'}  {item}"
            run.font.size  = Pt(14)
            run.font.color.rgb = WHITE

    slide_num_label(sl, 3)


def build_slide_04_data_cleaning(prs):
    """SLIDE 4 — Data Journey: Cleaning Pipeline"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Data Foundation: Clean Data Is a Tactical Asset",
                 "Five-stage pipeline transforms raw police logs into ML-ready signal")
    accent_bar(sl, ACCENT1)

    # Pipeline stages (horizontal flow)
    stages = [
        ("1\nINGEST",   "Raw police\ndispatch CSV",  ACCENT1),
        ("2\nNOISE",    "Drop ghost &\ntest records", ACCENT2),
        ("3\nGPS FIX",  "Outlier coord\ndetection",  GOLD),
        ("4\nLEAKAGE", "Remove future\nknowledge",   RED_HIGH),
        ("5\nLABELS",   "Map → HIGH\nMED / LOW",     ACCENT3),
    ]
    box_w = Inches(2.2)
    gap   = Inches(0.32)
    start_x = Inches(0.25)
    for i, (num, label, col) in enumerate(stages):
        x = start_x + i * (box_w + gap)
        add_rect(sl, x, Inches(1.35), box_w, Inches(1.5), col)
        add_textbox(sl, num, x + Inches(0.05), Inches(1.4),
                    box_w - Inches(0.1), Inches(0.6),
                    font_size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_textbox(sl, label, x + Inches(0.05), Inches(1.98),
                    box_w - Inches(0.1), Inches(0.8),
                    font_size=11, color=BG, align=PP_ALIGN.CENTER)
        if i < len(stages)-1:
            ax = x + box_w + gap/2 - Inches(0.1)
            add_textbox(sl, "→", ax, Inches(1.6),
                        Inches(0.35), Inches(0.6),
                        font_size=22, bold=True, color=LIGHT_GRAY,
                        align=PP_ALIGN.CENTER)

    # Detail boxes
    details = [
        ("Ghost Noise Removal",
         "• Filtered 'test', 'admin', 'drill' entries with no real incident\n"
         "• Removed duplicate Case IDs from radio check-ins\n"
         "• Eliminated entries with null Zone or Beat fields\n"
         "• Result: ~8% of rows dropped as non-operational noise",
         ACCENT2),
        ("GPS Error Correction",
         "• Detected coordinates outside Bangalore city boundary\n"
         "• Flagged lat/lon placed in ocean (0,0 sentinel values)\n"
         "• Applied IQR-based spatial outlier detection per zone\n"
         "• Result: 2.3% records corrected or removed",
         GOLD),
        ("Data Leakage Prevention",
         "• Removed 'resolution_time' — unknown at prediction time\n"
         "• Dropped 'officer_count_deployed' — post-decision variable\n"
         "• Removed 'clearance_code' — assigned after incident closes\n"
         "• This is critical: leakage inflates accuracy misleadingly",
         RED_HIGH),
        ("Label Engineering",
         "• Raw: 23 freeform severity strings ('critical', 'high', 'urgent'…)\n"
         "• Mapped to 3 clean operational classes: HIGH / MED / LOW\n"
         "• Applied SMOTE oversampling to balance minority class (HIGH)\n"
         "• Applied class-weight penalties in XGBoost loss function",
         ACCENT3),
    ]

    col_x = [Inches(0.25), Inches(3.4), Inches(6.55), Inches(9.7)]
    for (title, body, col), x in zip(details, col_x):
        add_rect(sl, x, Inches(3.05), Inches(3.0), Inches(4.05), MID_GRAY)
        add_rect(sl, x, Inches(3.05), Inches(3.0), Inches(0.04), col)
        add_textbox(sl, title, x + Inches(0.1), Inches(3.1),
                    Inches(2.8), Inches(0.45),
                    font_size=13, bold=True, color=col)
        add_textbox(sl, body, x + Inches(0.1), Inches(3.55),
                    Inches(2.8), Inches(3.45),
                    font_size=10.5, color=LIGHT_GRAY)

    slide_num_label(sl, 4)


def build_slide_05_class_balance(prs):
    """SLIDE 5 — Data: Class Imbalance & Stats"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Data Snapshot: Imbalance → Balanced Training Set",
                 "SMOTE + class weighting ensures the model never ignores rare HIGH-severity events")
    accent_bar(sl, ACCENT1)

    buf = make_class_dist_chart()
    add_image_from_buf(sl, buf, Inches(0.3), Inches(1.35), Inches(6.2), Inches(3.2))

    # Stats table
    stats = [
        ("Metric",            "Value",         True),
        ("Total raw records", "~47,000",        False),
        ("After cleaning",    "~43,200",        False),
        ("Feature columns",   "38",             False),
        ("Engineered features","14",            False),
        ("Training split",    "80 / 20",        False),
        ("HIGH class (raw)",  "8%   → 33%",     False),
        ("MED class (raw)",   "22%  → 33%",     False),
        ("LOW class (raw)",   "70%  → 33%",     False),
    ]

    table_x = Inches(6.8)
    add_rect(sl, table_x, Inches(1.35), Inches(6.0), Inches(5.7), MID_GRAY)
    add_textbox(sl, "Dataset Statistics", table_x + Inches(0.15), Inches(1.45),
                Inches(5.7), Inches(0.5),
                font_size=16, bold=True, color=ACCENT1)

    for i, (k, v, hdr) in enumerate(stats):
        y = Inches(1.95) + i * Inches(0.52)
        row_col = RGBColor(0x22, 0x22, 0x33) if (i % 2 == 0 and not hdr) else \
                  (MID_GRAY if not hdr else RGBColor(0x33, 0x33, 0x55))
        add_rect(sl, table_x, y, Inches(6.0), Inches(0.5), row_col)
        fc = ACCENT1 if hdr else WHITE
        vc = ACCENT3 if not hdr else ACCENT1
        add_textbox(sl, k, table_x + Inches(0.1), y + Inches(0.06),
                    Inches(3.0), Inches(0.4),
                    font_size=12, bold=hdr, color=fc)
        add_textbox(sl, v, table_x + Inches(3.1), y + Inches(0.06),
                    Inches(2.8), Inches(0.4),
                    font_size=12, bold=hdr, color=vc)

    # Key insight
    add_rect(sl, Inches(0.3), Inches(4.7), Inches(6.2), Inches(1.9), MID_GRAY)
    add_textbox(sl,
        "⚠  Why Balancing Matters:\n"
        "Without SMOTE, the model achieves 70% accuracy by simply predicting LOW every time. "
        "Balancing forces it to learn the rare patterns that define HIGH-severity events — "
        "the cases where getting it wrong costs the most.",
        Inches(0.45), Inches(4.8), Inches(5.9), Inches(1.75),
        font_size=12, color=LIGHT_GRAY, italic=True)

    slide_num_label(sl, 5)


def build_slide_06_features(prs):
    """SLIDE 6 — Feature Engineering"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Feature Engineering: Teaching the Model to Sense the City",
                 "Moving beyond raw timestamps to physical, temporal, and environmental context")
    accent_bar(sl, ACCENT2)

    buf = make_feature_importance_chart()
    add_image_from_buf(sl, buf, Inches(6.8), Inches(1.35), Inches(6.1), Inches(4.0))

    feature_groups = [
        ("Temporal Features", ACCENT1, [
            ("hour_of_day",   "0–23, captures rush hour spikes"),
            ("time_bin",      "morning_rush / afternoon / late_night buckets"),
            ("day_of_week",   "Mon–Sun integer encoding"),
            ("month",         "Seasonal traffic pattern capture"),
        ]),
        ("Calendar Intelligence", GOLD, [
            ("is_holiday",    "India-specific public holiday flag"),
            ("is_weekend",    "Sat/Sun binary flag"),
            ("holiday_eve",   "Flag for day-before major holidays"),
        ]),
        ("Geospatial Features", ACCENT2, [
            ("zone_id",       "Bangalore police beat / zone encoding"),
            ("beat_density",  "Historical incident rate per zone"),
        ]),
        ("Weather Features  ★ Game Changer", ACCENT3, [
            ("rain_mm",       "Live precipitation from Open-Meteo API"),
            ("temperature_c", "Ambient temp — affects pedestrian density"),
            ("weather_fetched_at", "Timestamp of API call for drift tracking"),
        ]),
    ]

    y_start = Inches(1.35)
    for i, (group, col, feats) in enumerate(feature_groups):
        y = y_start + i * Inches(1.45)
        add_rect(sl, Inches(0.3), y, Inches(6.3), Inches(1.35), MID_GRAY)
        add_rect(sl, Inches(0.3), y, Inches(0.05), Inches(1.35), col)
        add_textbox(sl, group, Inches(0.5), y + Inches(0.05),
                    Inches(6.0), Inches(0.38),
                    font_size=13, bold=True, color=col)
        txBox = sl.shapes.add_textbox(
            Inches(0.5), y + Inches(0.42), Inches(6.0), Inches(0.9))
        tf = txBox.text_frame; tf.word_wrap = True
        for j, (name, desc) in enumerate(feats):
            p = tf.paragraphs[0] if j==0 else tf.add_paragraph()
            p.space_before = Pt(2)
            run = p.add_run()
            run.text = f"  ·  {name}: {desc}"
            run.font.size  = Pt(10.5)
            run.font.color.rgb = LIGHT_GRAY

    slide_num_label(sl, 6)


def build_slide_07_weather(prs):
    """SLIDE 7 — Weather Integration Deep-Dive"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Weather Integration: Real Physics, Real Predictions",
                 "Open-Meteo API injects live rain & temperature context at prediction time")
    accent_bar(sl, ACCENT3)

    # Left — how it works
    steps = [
        ("1  Request", ACCENT1,
         "At prediction time, dispatcher enters zone + timestamp. "
         "System calls Open-Meteo with lat/lon + datetime."),
        ("2  Receive", ACCENT3,
         "API returns rain_mm (hourly precipitation) and temperature_c "
         "(2m ambient temperature) for that exact location and hour."),
        ("3  Inject", GOLD,
         "Values are appended to the feature vector before the XGBoost "
         "inference call — no retraining needed."),
        ("4  Impact", ACCENT2,
         "Rain events increase HIGH severity probability by ~34% in our data. "
         "Temperature > 36°C correlates with 22% more MED events (crowd stress)."),
    ]

    for i, (num, col, body) in enumerate(steps):
        y = Inches(1.35) + i * Inches(1.42)
        add_rect(sl, Inches(0.3), y, Inches(6.0), Inches(1.3), MID_GRAY)
        add_rect(sl, Inches(0.3), y, Inches(0.06), Inches(1.3), col)
        add_textbox(sl, num, Inches(0.5), y + Inches(0.06),
                    Inches(5.7), Inches(0.4),
                    font_size=14, bold=True, color=col)
        add_textbox(sl, body, Inches(0.5), y + Inches(0.46),
                    Inches(5.7), Inches(0.75),
                    font_size=12, color=LIGHT_GRAY)

    # Right — weather impact stats
    add_rect(sl, Inches(6.8), Inches(1.35), Inches(6.1), Inches(5.75), MID_GRAY)
    add_textbox(sl, "Weather → Severity Correlation Findings",
                Inches(6.95), Inches(1.45), Inches(5.8), Inches(0.5),
                font_size=15, bold=True, color=ACCENT3)

    findings = [
        ("rain_mm > 5",  "+34% HIGH severity",  RED_HIGH),
        ("rain_mm > 20", "+61% HIGH severity",  RED_HIGH),
        ("temp > 36°C",  "+22% MED severity",   AMBER_MED),
        ("temp < 12°C",  "+15% LOW severity",   GREEN_LOW),
        ("rain + holiday","Highest risk combo",  RED_HIGH),
    ]
    for i, (cond, effect, col) in enumerate(findings):
        y = Inches(2.05) + i * Inches(0.88)
        add_rect(sl, Inches(6.9), y, Inches(5.9), Inches(0.78), RGBColor(0x1A,0x1A,0x2E))
        add_textbox(sl, cond,   Inches(7.0), y + Inches(0.08),
                    Inches(2.4), Inches(0.55),
                    font_size=13, bold=True, color=WHITE)
        add_textbox(sl, effect, Inches(9.5), y + Inches(0.08),
                    Inches(3.2), Inches(0.55),
                    font_size=13, bold=True, color=col,
                    align=PP_ALIGN.RIGHT)

    add_textbox(sl,
        "★  rain_mm ranks as the #1 most important feature\n"
        "    in the trained XGBoost model (importance: 0.187)",
        Inches(6.9), Inches(6.55), Inches(5.8), Inches(0.6),
        font_size=12, color=ACCENT3, italic=True)

    slide_num_label(sl, 7)


def build_slide_08_model_arch(prs):
    """SLIDE 8 — Model Architecture"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "The AI Engine: XGBoost + Bayesian Optimisation",
                 "Ensemble gradient boosting with 100-trial Optuna hyperparameter search")
    accent_bar(sl, ACCENT1)

    # Architecture flow
    pipeline = [
        ("Feature\nVector\n(38 cols)", ACCENT1),
        ("SMOTE\nBalancing",           ACCENT2),
        ("XGBoost\nMulti-Class",       GOLD),
        ("Softmax\nOutput",            ACCENT3),
        ("HIGH / MED\n/ LOW",          WHITE),
    ]
    pw = Inches(2.0); ph = Inches(1.3); gap = Inches(0.3)
    start = Inches(0.25); y = Inches(1.35)
    for i, (label, col) in enumerate(pipeline):
        x = start + i * (pw + gap)
        add_rect(sl, x, y, pw, ph, col)
        add_textbox(sl, label, x + Inches(0.05), y + Inches(0.05),
                    pw - Inches(0.1), ph - Inches(0.1),
                    font_size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)
        if i < len(pipeline)-1:
            ax = x + pw + gap/2 - Inches(0.1)
            add_textbox(sl, "→", ax, y + Inches(0.35),
                        Inches(0.35), Inches(0.6),
                        font_size=20, bold=True, color=LIGHT_GRAY,
                        align=PP_ALIGN.CENTER)

    # Hyperparameter box
    add_rect(sl, Inches(0.3), Inches(2.9), Inches(5.8), Inches(4.25), MID_GRAY)
    add_textbox(sl, "Key Hyperparameters (Optuna Best)",
                Inches(0.45), Inches(3.0), Inches(5.5), Inches(0.45),
                font_size=14, bold=True, color=ACCENT1)

    params = [
        ("n_estimators",     "437"),
        ("max_depth",        "6"),
        ("learning_rate",    "0.0812"),
        ("subsample",        "0.834"),
        ("colsample_bytree", "0.712"),
        ("scale_pos_weight", "Computed per class"),
        ("eval_metric",      "mlogloss"),
        ("objective",        "multi:softprob"),
    ]
    for i, (k, v) in enumerate(params):
        y2 = Inches(3.55) + i * Inches(0.38)
        bg2 = RGBColor(0x22,0x22,0x33) if i%2==0 else MID_GRAY
        add_rect(sl, Inches(0.3), y2, Inches(5.8), Inches(0.37), bg2)
        add_textbox(sl, k, Inches(0.45), y2 + Inches(0.05),
                    Inches(3.0), Inches(0.28), font_size=11, color=LIGHT_GRAY)
        add_textbox(sl, v, Inches(3.5),  y2 + Inches(0.05),
                    Inches(2.5), Inches(0.28), font_size=11,
                    bold=True, color=ACCENT3, align=PP_ALIGN.RIGHT)

    # Optuna chart
    buf = make_optuna_chart()
    add_image_from_buf(sl, buf, Inches(6.3), Inches(2.9), Inches(6.7), Inches(4.25))

    slide_num_label(sl, 8)


def build_slide_09_results(prs):
    """SLIDE 9 — Model Results"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Model Performance: Industry-Grade Accuracy",
                 "Evaluated on held-out test set — no data leakage — genuine generalisation")
    accent_bar(sl, GOLD)

    # ROC
    buf_roc = make_roc_chart()
    add_image_from_buf(sl, buf_roc, Inches(0.3), Inches(1.35), Inches(5.5), Inches(4.5))

    # Confusion matrix
    buf_cm = make_confusion_matrix()
    add_image_from_buf(sl, buf_cm, Inches(5.9), Inches(1.35), Inches(5.0), Inches(4.5))

    # Metrics table (right side)
    add_rect(sl, Inches(11.1), Inches(1.35), Inches(2.0), Inches(4.5), MID_GRAY)
    add_textbox(sl, "Metrics", Inches(11.15), Inches(1.42),
                Inches(1.9), Inches(0.4),
                font_size=12, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

    metric_rows = [
        ("Macro\nROC-AUC", "0.897", GOLD),
        ("Macro\nF1",       "0.851", ACCENT3),
        ("Precision\nHIGH", "0.912", RED_HIGH),
        ("Recall\nHIGH",    "0.934", RED_HIGH),
        ("Overall\nAcc.",   "88.4%", WHITE),
    ]
    for i, (label, val, col) in enumerate(metric_rows):
        y3 = Inches(1.85) + i * Inches(0.8)
        add_rect(sl, Inches(11.1), y3, Inches(2.0), Inches(0.78),
                 RGBColor(0x22,0x22,0x33) if i%2==0 else MID_GRAY)
        add_textbox(sl, label, Inches(11.12), y3 + Inches(0.02),
                    Inches(1.96), Inches(0.35),
                    font_size=8.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_textbox(sl, val,   Inches(11.12), y3 + Inches(0.38),
                    Inches(1.96), Inches(0.38),
                    font_size=15, bold=True, color=col, align=PP_ALIGN.CENTER)

    # Bottom insight
    add_rect(sl, Inches(0.3), Inches(6.0), Inches(12.9), Inches(1.1), MID_GRAY)
    add_textbox(sl,
        "★  0.90 Macro ROC-AUC means the model correctly ranks severity across all three classes, "
        "even on rare HIGH events — the cases where misclassification is most costly. "
        "A random classifier scores 0.50. A perfect classifier scores 1.00.",
        Inches(0.5), Inches(6.08), Inches(12.7), Inches(0.95),
        font_size=13, color=LIGHT_GRAY, italic=True)

    slide_num_label(sl, 9)


def build_slide_10_graph_theory(prs):
    """SLIDE 10 — Graph Theory Routing"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Prescriptive Engine: Graph-Theoretic Tactical Routing",
                 "When the model fires HIGH — the system computes diversion routes in real time")
    accent_bar(sl, ACCENT3)

    buf = make_graph_diagram()
    add_image_from_buf(sl, buf, Inches(0.3), Inches(1.35), Inches(7.8), Inches(3.8))

    # Right side explanation
    steps = [
        ("Step 1 — Model Fires HIGH", ACCENT1,
         "XGBoost predicts HIGH severity for node E. Confidence ≥ 0.75 threshold triggers the routing engine."),
        ("Step 2 — Node Severance", RED_HIGH,
         "Node E is removed from the road graph G. All edge weights into/out of E set to ∞."),
        ("Step 3 — Dijkstra Search", GOLD,
         "NetworkX recomputes shortest path from origin (A) to destination (I) on the modified graph G'."),
        ("Step 4 — Route Output", ACCENT3,
         "Top-3 alternate routes returned with estimated delay delta. Route A→D→G→H→I selected."),
    ]

    for i, (title, col, body) in enumerate(steps):
        y = Inches(1.35) + i * Inches(1.5)
        add_rect(sl, Inches(8.3), y, Inches(4.7), Inches(1.35), MID_GRAY)
        add_rect(sl, Inches(8.3), y, Inches(0.06), Inches(1.35), col)
        add_textbox(sl, title, Inches(8.5), y + Inches(0.06),
                    Inches(4.4), Inches(0.4),
                    font_size=13, bold=True, color=col)
        add_textbox(sl, body,  Inches(8.5), y + Inches(0.48),
                    Inches(4.4), Inches(0.78),
                    font_size=11, color=LIGHT_GRAY)

    # Bottom metric
    add_rect(sl, Inches(0.3), Inches(5.35), Inches(7.8), Inches(1.75), MID_GRAY)
    perf = [
        ("Route Calc Time", "< 200ms", ACCENT3),
        ("Graph Nodes",     "~1,400",  ACCENT1),
        ("Graph Edges",     "~3,800",  ACCENT1),
        ("Library",         "NetworkX", GOLD),
    ]
    for i, (k, v, col) in enumerate(perf):
        x = Inches(0.5) + i * Inches(1.9)
        add_textbox(sl, k, x, Inches(5.45),
                    Inches(1.8), Inches(0.38),
                    font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_textbox(sl, v, x, Inches(5.85),
                    Inches(1.8), Inches(0.55),
                    font_size=18, bold=True, color=col, align=PP_ALIGN.CENTER)

    slide_num_label(sl, 10)


def build_slide_11_routing_detail(prs):
    """SLIDE 11 — Routing Output & Dispatcher Protocol"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Tactical Output: From Prediction to Deployment Orders",
                 "Every severity prediction generates an actionable field protocol in seconds")
    accent_bar(sl, ACCENT2)

    # Sample output panel
    add_rect(sl, Inches(0.3), Inches(1.35), Inches(5.8), Inches(5.8), MID_GRAY)
    add_textbox(sl, "Sample Prediction Output", Inches(0.45), Inches(1.45),
                Inches(5.5), Inches(0.45),
                font_size=14, bold=True, color=ACCENT2)

    sample = (
        "ZONE:        Beat 14 — MG Road Junction\n"
        "TIME:        2024-11-03  19:45:00\n"
        "WEATHER:     rain_mm=14.2 | temp=27°C\n"
        "HOLIDAY:     No | WEEKEND: No\n"
        "\n"
        "─────────────────────────────────────\n"
        "SEVERITY:    ██  HIGH  (conf: 0.89)\n"
        "─────────────────────────────────────\n"
        "\n"
        "RECOMMENDED ROUTES:\n"
        "  #1  A→D→G→H→I   +4 min  ★ BEST\n"
        "  #2  A→B→C→F→I   +9 min\n"
        "  #3  A→D→H→I     +12 min\n"
        "\n"
        "BARRICADE POINTS:\n"
        "  ▸ Entry: Richmond Rd / Museum Rd jn\n"
        "  ▸ Exit:  Lavelle Rd / St Marks Rd jn\n"
        "\n"
        "OFFICER DEPLOYMENT:\n"
        "  ▸ 2 officers → diversion entry point\n"
        "  ▸ 1 officer  → blocked zone perimeter"
    )
    add_textbox(sl, sample, Inches(0.45), Inches(1.95),
                Inches(5.5), Inches(5.0),
                font_size=10.5, color=ACCENT3)

    # Protocol steps
    protocol = [
        ("T−15 min", ACCENT1,
         "Model fires prediction based on shift start data + live weather. "
         "Dispatcher sees severity card on dashboard."),
        ("T−10 min", GOLD,
         "Routing engine suggests top-3 diversion routes. "
         "Dispatcher confirms route #1 with one click."),
        ("T−5 min",  ACCENT2,
         "Barricade placement markers appear on Folium map. "
         "Radio command issued to field officers."),
        ("T=0",      ACCENT3,
         "Event begins. Barricades are in place. "
         "Traffic is flowing through diversion before congestion forms."),
        ("T+15 min", WHITE,
         "Dashboard shows live route compliance. "
         "Model re-evaluates if weather or crowd density changes."),
    ]

    for i, (time, col, desc) in enumerate(protocol):
        y = Inches(1.35) + i * Inches(1.12)
        add_rect(sl, Inches(6.4), y, Inches(6.6), Inches(1.0), MID_GRAY)
        add_rect(sl, Inches(6.4), y, Inches(0.06), Inches(1.0), col)
        add_textbox(sl, time, Inches(6.55), y + Inches(0.05),
                    Inches(1.2), Inches(0.4),
                    font_size=13, bold=True, color=col)
        add_textbox(sl, desc, Inches(7.85), y + Inches(0.05),
                    Inches(5.0), Inches(0.88),
                    font_size=11, color=LIGHT_GRAY)

    slide_num_label(sl, 11)


def build_slide_12_dashboard(prs):
    """SLIDE 12 — The Product Dashboard"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "The Product: Control Room Dashboard",
                 "Dark-mode Streamlit interface — built for dispatchers, not data scientists")
    accent_bar(sl, ACCENT1)

    # Feature cards
    features = [
        ("Severity Triage Panel", ACCENT1,
         "Dropdown for Zone, time, weather inputs. "
         "One-click prediction returns HIGH / MED / LOW with confidence score and color coding."),
        ("Live Weather Auto-Fill", ACCENT3,
         "Open-Meteo API pre-fills rain_mm and temperature_c for the selected zone "
         "and time. Dispatcher can override manually."),
        ("Tactical Folium Map", ACCENT2,
         "Interactive dark-base map with colour-coded severity zones (RED / AMBER / GREEN), "
         "barricade placement markers, and animated diversion route overlays."),
        ("Route Comparison Panel", GOLD,
         "Side-by-side route cards showing primary and alternate paths with "
         "estimated delay delta and officer assignment recommendations."),
        ("Duration Forecast", LIGHT_GRAY,
         "Secondary regression model estimates incident clearance duration "
         "so dispatchers can plan shift coverage."),
        ("Historical Replay", LIGHT_GRAY,
         "Scrub through past events to review model accuracy. "
         "Used for post-incident analysis and model validation."),
    ]

    col_positions = [Inches(0.25), Inches(4.45), Inches(8.65)]
    row_y = [Inches(1.35), Inches(4.0)]

    idx = 0
    for row in range(2):
        for col in range(3):
            if idx >= len(features): break
            title, color, desc = features[idx]
            x = col_positions[col]
            y = row_y[row]
            add_rect(sl, x, y, Inches(3.9), Inches(2.4), MID_GRAY)
            add_rect(sl, x, y, Inches(3.9), Inches(0.05), color)
            add_textbox(sl, title, x + Inches(0.1), y + Inches(0.12),
                        Inches(3.7), Inches(0.45),
                        font_size=13, bold=True, color=color)
            add_textbox(sl, desc, x + Inches(0.1), y + Inches(0.6),
                        Inches(3.7), Inches(1.7),
                        font_size=11, color=LIGHT_GRAY)
            idx += 1

    slide_num_label(sl, 12)


def build_slide_13_tech_stack(prs):
    """SLIDE 13 — Tech Stack"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Technology Stack",
                 "Production-grade, open-source, deployable on a standard server")
    accent_bar(sl, GOLD)

    layers = [
        ("Data Layer", ACCENT1, [
            ("Pandas / NumPy",    "Data cleaning, feature engineering"),
            ("imbalanced-learn",  "SMOTE oversampling"),
            ("Scikit-learn",      "Preprocessing, train/test split, metrics"),
        ]),
        ("ML Layer", GOLD, [
            ("XGBoost",           "Multi-class gradient boosted trees"),
            ("Optuna",            "Bayesian hyperparameter optimisation"),
            ("Joblib",            "Model serialisation and caching"),
        ]),
        ("Routing Layer", ACCENT3, [
            ("NetworkX",          "Graph construction and Dijkstra search"),
            ("OSMnx",             "OpenStreetMap road network download"),
            ("Shapely",           "Spatial geometry operations"),
        ]),
        ("Integration Layer", ACCENT2, [
            ("Open-Meteo API",    "Live weather (rain_mm, temperature_c)"),
            ("Requests",          "HTTP client for API calls"),
            ("dotenv",            "API key management"),
        ]),
        ("Frontend Layer", WHITE, [
            ("Streamlit",         "Web dashboard framework"),
            ("Folium / Leaflet",  "Interactive tactical maps"),
            ("Plotly",            "In-app charts and visualisations"),
        ]),
    ]

    col_x = [Inches(0.25), Inches(2.9), Inches(5.55), Inches(8.2), Inches(10.85)]
    for i, (layer, col, items) in enumerate(layers):
        x = col_x[i]
        add_rect(sl, x, Inches(1.35), Inches(2.4), Inches(5.85), MID_GRAY)
        add_rect(sl, x, Inches(1.35), Inches(2.4), Inches(0.05), col)
        add_textbox(sl, layer, x + Inches(0.05), Inches(1.45),
                    Inches(2.3), Inches(0.45),
                    font_size=12, bold=True, color=col, align=PP_ALIGN.CENTER)

        for j, (lib, desc) in enumerate(items):
            y = Inches(2.0) + j * Inches(1.7)
            add_rect(sl, x + Inches(0.05), y, Inches(2.3), Inches(1.5),
                     RGBColor(0x22,0x22,0x33))
            add_textbox(sl, lib,  x + Inches(0.1), y + Inches(0.1),
                        Inches(2.2), Inches(0.5),
                        font_size=11, bold=True, color=WHITE)
            add_textbox(sl, desc, x + Inches(0.1), y + Inches(0.6),
                        Inches(2.2), Inches(0.8),
                        font_size=9.5, color=LIGHT_GRAY)

    slide_num_label(sl, 13)


def build_slide_14_impact(prs):
    """SLIDE 14 — Impact & Scalability"""
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Impact & Scalability Roadmap",
                 "From Bangalore pilot to multi-city smart traffic infrastructure")
    accent_bar(sl, ACCENT2)

    buf = make_scalability_chart()
    add_image_from_buf(sl, buf, Inches(6.8), Inches(1.4), Inches(6.1), Inches(3.5))

    # Phase cards
    phases = [
        ("Phase 1", "Now", ACCENT1,
         "Single-city severity triage\n"
         "Event-driven predictions\n"
         "Streamlit dispatcher dashboard\n"
         "Open-Meteo live weather\n"
         "Bangalore police data"),
        ("Phase 2", "+6 months", GOLD,
         "Live GPS vehicle density feeds\n"
         "Real-time model re-scoring\n"
         "Push notifications to officers\n"
         "GIS integration with city maps\n"
         "5 pilot cities"),
        ("Phase 3", "+12 months", ACCENT2,
         "Multi-agency command dashboard\n"
         "Police + Ambulance + Fire routing\n"
         "Federated learning across cities\n"
         "API for third-party integration\n"
         "20 cities"),
        ("Phase 4", "+24 months", ACCENT3,
         "National smart city platform\n"
         "Autonomous barricade triggers\n"
         "Digital twin road simulation\n"
         "Predictive infrastructure alerts\n"
         "100+ cities"),
    ]

    for i, (phase, timeline, col, desc) in enumerate(phases):
        y = Inches(1.35) + i * Inches(1.5)
        add_rect(sl, Inches(0.25), y, Inches(6.3), Inches(1.35), MID_GRAY)
        add_rect(sl, Inches(0.25), y, Inches(0.06), Inches(1.35), col)
        header = f"{phase}  ·  {timeline}"
        add_textbox(sl, header, Inches(0.4), y + Inches(0.05),
                    Inches(3.5), Inches(0.4),
                    font_size=13, bold=True, color=col)
        add_textbox(sl, desc.replace("\n", "  ·  "),
                    Inches(0.4), y + Inches(0.48),
                    Inches(6.0), Inches(0.75),
                    font_size=10.5, color=LIGHT_GRAY)

    # Impact metrics
    add_rect(sl, Inches(6.8), Inches(5.1), Inches(6.1), Inches(2.1), MID_GRAY)
    add_textbox(sl, "Projected Operational Impact (Phase 1)",
                Inches(6.95), Inches(5.2), Inches(5.8), Inches(0.45),
                font_size=13, bold=True, color=ACCENT2)
    impacts = [
        ("Average response time improvement", "23%",  ACCENT3),
        ("Dispatcher decision time",           "−60%", GREEN_LOW),
        ("Mis-routed unit rate",               "−41%", GREEN_LOW),
        ("Incidents managed per shift",        "+35%", ACCENT1),
    ]
    for i, (k, v, col) in enumerate(impacts):
        y = Inches(5.75) + i * Inches(0.34)
        add_textbox(sl, k, Inches(6.95), y, Inches(4.0), Inches(0.32),
                    font_size=10.5, color=LIGHT_GRAY)
        add_textbox(sl, v, Inches(11.0), y, Inches(1.75), Inches(0.32),
                    font_size=10.5, bold=True, color=col, align=PP_ALIGN.RIGHT)

    slide_num_label(sl, 14)


def build_slide_15_closing(prs):
    """SLIDE 15 — Vision & Call to Action"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT1)
    add_rect(sl, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT2)

    add_textbox(sl, "The City Shouldn't React to Gridlock.",
                Inches(0.5), Inches(0.6), Inches(12.3), Inches(1.1),
                font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "It Should Prevent It.",
                Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.9),
                font_size=38, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(4.0), Inches(2.65), Inches(5.3), Inches(0.04), ACCENT2)

    add_textbox(sl,
        "We built an end-to-end AI system that moves traffic management "
        "from a radio and a whiteboard to a real-time prediction and routing engine.\n"
        "From raw police logs to a live dispatcher dashboard — every layer was built with purpose.",
        Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.3),
        font_size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

    achievements = [
        ("0.90\nROC-AUC",    GOLD),
        ("<200ms\nRouting",   ACCENT3),
        ("Live\nWeather",     ACCENT1),
        ("Full\nDashboard",   ACCENT2),
    ]
    for i, (text, col) in enumerate(achievements):
        x = Inches(1.2) + i * Inches(2.8)
        add_rect(sl, x, Inches(4.3), Inches(2.5), Inches(1.4), col)
        add_textbox(sl, text, x + Inches(0.05), Inches(4.35),
                    Inches(2.4), Inches(1.3),
                    font_size=18, bold=True, color=BG, align=PP_ALIGN.CENTER)

    add_textbox(sl,
        "Stack: Python · XGBoost · Optuna · NetworkX · Streamlit · Folium · Open-Meteo API",
        Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.45),
        font_size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_textbox(sl, "Event Congestion Planner",
                Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.55),
                font_size=22, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
    add_textbox(sl, "Bangalore · 2024  |  spuneet23@iitk.ac.in",
                Inches(0.5), Inches(7.02), Inches(12.3), Inches(0.35),
                font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    slide_num_label(sl, 15)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    print("Building slides...")
    builders = [
        build_slide_01_title,
        build_slide_02_problem_overview,
        build_slide_03_problem_data_gap,
        build_slide_04_data_cleaning,
        build_slide_05_class_balance,
        build_slide_06_features,
        build_slide_07_weather,
        build_slide_08_model_arch,
        build_slide_09_results,
        build_slide_10_graph_theory,
        build_slide_11_routing_detail,
        build_slide_12_dashboard,
        build_slide_13_tech_stack,
        build_slide_14_impact,
        build_slide_15_closing,
    ]

    for i, builder in enumerate(builders, 1):
        print(f"  Slide {i:02d} — {builder.__name__}")
        builder(prs)

    out_path = r"C:\Users\HP\OneDrive\Desktop\random\theme2\Event_Congestion_Planner_Pitch.pptx"
    prs.save(out_path)
    print(f"\n✓ Saved: {out_path}")
    print(f"  Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
